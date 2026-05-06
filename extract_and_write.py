import os
import pptx
from pptx import Presentation

pptx_path = r'C:\Users\Márcio\Desktop\Tornado_Building_Simulator\Final Presentation - 05 01 26 - Márcio Ferreira.pptx'
text_output_path = r'C:\Users\Márcio\Desktop\Tornado_Building_Simulator\extracted_text.txt'
images_dir = r'C:\Users\Márcio\Desktop\Tornado_Building_Simulator\extracted_images'

os.makedirs(images_dir, exist_ok=True)

prs = Presentation(pptx_path)

extracted_text = []
image_count = 0

for i, slide in enumerate(prs.slides):
    extracted_text.append(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            extracted_text.append(shape.text.strip())
            
        # Extract images
        if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
            if image_count < 15:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_filename = os.path.join(images_dir, f'image_{i+1}_{image_count}.{image_ext}')
                with open(image_filename, 'wb') as f:
                    f.write(image_bytes)
                extracted_text.append(f"[Image extracted: {image_filename}]")
                image_count += 1
                
        # Also check groups for images
        if shape.shape_type == 6: # Group
            for child in shape.shapes:
                if child.shape_type == 13 and image_count < 15:
                    image = child.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_filename = os.path.join(images_dir, f'image_{i+1}_group_{image_count}.{image_ext}')
                    with open(image_filename, 'wb') as f:
                        f.write(image_bytes)
                    extracted_text.append(f"[Image extracted: {image_filename}]")
                    image_count += 1

with open(text_output_path, 'w', encoding='utf-8') as f:
    f.write("\n".join(extracted_text))

print(f"Extracted text saved to {text_output_path}")
print(f"Extracted {image_count} images to {images_dir}")
